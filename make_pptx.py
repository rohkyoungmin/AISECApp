"""Reporter 중간 발표 PPTX 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── 색상 팔레트 (Reporter 디자인 시스템) ──────────────────────────────────────
C_PRIMARY   = RGBColor(0x01, 0x40, 0x29)   # 다크 그린 #014029
C_CANVAS    = RGBColor(0xFA, 0xF9, 0xF5)   # 크림 배경
C_CARD      = RGBColor(0xEF, 0xE9, 0xDE)   # 카드
C_DARK      = RGBColor(0x18, 0x17, 0x15)   # 다크
C_INK       = RGBColor(0x14, 0x14, 0x13)   # 본문 텍스트
C_MUTED     = RGBColor(0x6C, 0x6A, 0x64)   # 보조 텍스트
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_ERROR     = RGBColor(0xC6, 0x45, 0x45)   # Critical
C_AMBER     = RGBColor(0xE8, 0xA5, 0x5A)   # High
C_WARNING   = RGBColor(0xD4, 0xA0, 0x17)   # Medium
C_TEAL      = RGBColor(0x5D, 0xB8, 0xA6)   # Low
C_SUCCESS   = RGBColor(0x1A, 0x7A, 0x45)   # 성공
C_HAIRLINE  = RGBColor(0xE6, 0xDF, 0xD8)   # 구분선
C_ACCENT    = RGBColor(0xB0, 0xD8, 0xC0)   # 연그린 (헤더 서브텍스트)

W, H = Inches(13.33), Inches(7.5)  # 와이드스크린 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # 완전 빈 레이아웃


def add_rect(slide, x, y, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def add_multiline(slide, lines, x, y, w, h, size=14, color=None, bold=False, line_space=1.2):
    """lines = list of (text, bold_override) tuples or plain strings"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, tuple):
            txt, b = item
        else:
            txt, b = item, bold
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = b
        if color:
            run.font.color.rgb = color
    return tb


def dark_header(slide, title, subtitle=None):
    """공통: 다크그린 전체 배경 + 흰 제목"""
    add_rect(slide, 0, 0, 13.33, 7.5, C_PRIMARY)
    add_text(slide, title, 0.7, 2.5, 12, 1.2,
             size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.7, 3.9, 12, 0.7,
                 size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)


def section_slide(slide, section_num, title, subtitle=None):
    """섹션 구분 슬라이드"""
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK)
    add_rect(slide, 0, 0, 0.12, 7.5, C_PRIMARY)
    add_text(slide, section_num, 0.5, 1.8, 2, 1.2,
             size=72, bold=True, color=C_PRIMARY)
    add_text(slide, title, 0.5, 3.2, 12, 1.2,
             size=40, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 4.6, 11, 0.7,
                 size=18, color=C_ACCENT)


def content_slide(slide, title):
    """공통: 크림 배경 + 다크그린 상단 바"""
    add_rect(slide, 0, 0, 13.33, 7.5, C_CANVAS)
    add_rect(slide, 0, 0, 13.33, 1.05, C_PRIMARY)
    add_rect(slide, 0, 1.05, 13.33, 0.04, C_HAIRLINE)
    add_text(slide, title, 0.5, 0.15, 12, 0.75,
             size=28, bold=True, color=C_WHITE)


def card(slide, x, y, w, h, title=None, title_color=None):
    add_rect(slide, x, y, w, h, C_CARD)
    if title:
        add_rect(slide, x, y, w, 0.05, title_color or C_PRIMARY)
        add_text(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.4,
                 size=13, bold=True, color=C_INK)


def pill(slide, x, y, label, color):
    """작은 배지"""
    add_rect(slide, x, y, 1.2, 0.32, color)
    add_text(slide, label, x, y + 0.04, 1.2, 0.28,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, C_DARK)
add_rect(sl, 0, 0, 0.18, 7.5, C_PRIMARY)          # 왼쪽 녹색 스트라이프
add_rect(sl, 0, 6.2, 13.33, 1.3, C_PRIMARY)       # 하단 바

# 메인 타이틀
add_text(sl, "Reporter", 0.6, 1.3, 12, 1.8,
         size=82, bold=True, color=C_WHITE)
add_text(sl, "AI 기반 C/C++ 소스코드 취약점 분석 플랫폼", 0.6, 3.3, 12, 0.8,
         size=24, color=C_ACCENT)
add_text(sl, "Claude AI · 멀티 에이전트 파이프라인 · 자동 PDF 리포트", 0.6, 4.1, 12, 0.6,
         size=16, color=C_MUTED)

# 하단 바 정보
add_text(sl, "인공지능보안응용  ·  중간 발표  ·  2025.05", 0.6, 6.35, 8, 0.5,
         size=14, color=C_ACCENT)
add_text(sl, "노경민", 10.5, 6.35, 2.5, 0.5,
         size=14, color=C_WHITE, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 2 — 프로젝트 한 줄 소개
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "Reporter — 무엇을 만들었나")

add_text(sl, "ZIP 아카이브를 업로드하면, AI가 C/C++ 코드를 분석해 취약점을 찾고\n증거와 함께 PDF 리포트로 자동 출력합니다.",
         0.7, 1.3, 12, 1.0, size=20, color=C_INK)

# 4단계 파이프라인 흐름도
stages = [
    ("01\nExtract",  "ZIP 해제\nC/C++ 필터링", C_PRIMARY),
    ("02\nTriage",   "위험 패턴\n휴리스틱 탐지", C_SUCCESS),
    ("03\nFinding",  "Claude AI\n취약점 심층 분석", C_AMBER),
    ("04\nVerify\n→ Report", "증거 검증\nPDF 생성", C_ERROR),
]
for i, (num, desc, col) in enumerate(stages):
    x = 0.7 + i * 3.1
    add_rect(sl, x, 2.5, 2.8, 2.9, C_CARD)
    add_rect(sl, x, 2.5, 2.8, 0.06, col)
    add_text(sl, num, x + 0.15, 2.6, 2.5, 0.9,
             size=15, bold=True, color=col)
    add_text(sl, desc, x + 0.15, 3.55, 2.5, 0.9,
             size=13, color=C_INK)
    # 화살표
    if i < 3:
        add_text(sl, "→", x + 2.85, 3.7, 0.35, 0.4,
                 size=20, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

add_text(sl, "FastAPI 백엔드  ·  React + Vite 프론트엔드  ·  fpdf2 PDF 엔진  ·  Anthropic Claude Sonnet 4.6",
         0.7, 5.7, 12, 0.5, size=13, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 3 — 진행 상황 요약
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "진행 상황 요약 — 목표 대비 현재")

# 왼쪽: 완료
add_rect(sl, 0.5, 1.25, 5.9, 5.8, C_CARD)
add_rect(sl, 0.5, 1.25, 5.9, 0.06, C_SUCCESS)
add_text(sl, "✅  완료된 기능", 0.7, 1.35, 5.5, 0.45, size=15, bold=True, color=C_SUCCESS)

done = [
    "멀티 에이전트 분석 파이프라인 구현",
    "Claude API 연동 (Finding Agent)",
    "HeuristicTriageAgent — 파일당 API 1회로 최적화",
    "EvidencePolicyVerifier — 증거 기반 결정론적 검증",
    "Exploitability Policy (severity별 차등 기준)",
    "FastAPI 백엔드 + REST API 전체",
    "React + Vite 프론트엔드 (6개 페이지)",
    "워밍 크림 디자인 시스템 전면 적용",
    "PDF 리포트 자동 생성 (DejaVu 폰트)",
    "GitHub 저장소 관리 및 문서화",
]
for i, item in enumerate(done):
    add_text(sl, f"•  {item}", 0.7, 1.95 + i * 0.47, 5.5, 0.45,
             size=12, color=C_INK)

# 오른쪽: 남은 것
add_rect(sl, 6.9, 1.25, 5.9, 5.8, C_CARD)
add_rect(sl, 6.9, 1.25, 5.9, 0.06, C_AMBER)
add_text(sl, "🔲  남은 과제 (5/30까지)", 7.1, 1.35, 5.5, 0.45, size=15, bold=True, color=C_AMBER)

todo = [
    "CVE 케이스셋 정량 평가 (Precision / Recall)",
    "Confidence 보정 및 False Positive 분석",
    "프론트엔드 실시간 로그 스트리밍 (SSE)",
    "대용량 ZIP 처리 성능 최적화",
    "최종 Demo Day 시연 준비",
]
for i, item in enumerate(todo):
    add_text(sl, f"•  {item}", 7.1, 1.95 + i * 0.6, 5.5, 0.55,
             size=12, color=C_INK)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 4 — 시스템 아키텍처
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "시스템 아키텍처")

# 세 레이어
layers = [
    ("사용자  /  브라우저", "React + Vite  |  포트 5173", C_PRIMARY, 1.2),
    ("백엔드  API", "FastAPI  |  포트 8000  |  /projects  /jobs  /reports", C_SUCCESS, 2.6),
    ("분석 파이프라인", "MultiAgentSourceAnalyzer  →  ClaudeClient  →  Anthropic API", C_AMBER, 4.0),
]
for label, detail, col, y in layers:
    add_rect(sl, 1.2, y, 11.0, 0.9, C_CARD)
    add_rect(sl, 1.2, y, 0.07, 0.9, col)
    add_text(sl, label, 1.5, y + 0.08, 3.5, 0.35, size=14, bold=True, color=C_INK)
    add_text(sl, detail, 5.0, y + 0.08, 7.0, 0.35, size=12, color=C_MUTED)
    if y < 4.0:
        add_text(sl, "↕", 6.5, y + 0.95, 0.5, 0.35,
                 size=16, color=C_MUTED, align=PP_ALIGN.CENTER)

# Vite proxy
add_text(sl, "Vite Dev Proxy\n/projects, /jobs, /health  →  localhost:8000",
         3.5, 5.3, 7.0, 0.8, size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

add_text(sl, "ZIP 업로드 → 분석 Job 생성 → 파일별 파이프라인 실행 → Report 저장 → PDF 다운로드",
         0.7, 6.4, 12, 0.5, size=12, color=C_PRIMARY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 5 — 파이프라인 상세 (API 호출 최적화 전/후)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "AI 에이전트 파이프라인 — Claude API 호출 최적화")

# Before
add_rect(sl, 0.5, 1.25, 5.7, 5.8, C_CARD)
add_rect(sl, 0.5, 1.25, 5.7, 0.06, C_ERROR)
add_text(sl, "🔴  최적화 전  (파일당 API 3회)", 0.7, 1.35, 5.3, 0.45,
         size=14, bold=True, color=C_ERROR)

before = [
    ("ClaudeTriageAgent",    "API call ①  위험 패턴 판단", C_ERROR),
    ("ClaudeFindingAgent",   "API call ②  취약점 심층 분석", C_ERROR),
    ("ClaudeSkepticVerifier","API call ③  증거 재검증", C_ERROR),
    ("SourceReporterAgent",  "리포트 조립", C_MUTED),
]
for i, (name, desc, col) in enumerate(before):
    y = 2.0 + i * 1.1
    add_rect(sl, 0.7, y, 5.1, 0.85, C_CANVAS)
    add_rect(sl, 0.7, y, 0.06, 0.85, col)
    add_text(sl, name, 0.95, y + 0.07, 4.6, 0.32, size=12, bold=True, color=C_INK)
    add_text(sl, desc, 0.95, y + 0.42, 4.6, 0.32, size=11, color=C_MUTED)

# After
add_rect(sl, 7.0, 1.25, 5.8, 5.8, C_CARD)
add_rect(sl, 7.0, 1.25, 5.8, 0.06, C_SUCCESS)
add_text(sl, "✅  최적화 후  (파일당 API 1회)", 7.2, 1.35, 5.4, 0.45,
         size=14, bold=True, color=C_SUCCESS)

after = [
    ("HeuristicTriageAgent",   "정규식 기반 · API 없음 ✓", C_SUCCESS),
    ("ClaudeFindingAgent",     "API call ①  취약점 심층 분석", C_AMBER),
    ("EvidencePolicyVerifier", "결정론적 검증 · API 없음 ✓", C_SUCCESS),
    ("SourceReporterAgent",    "리포트 조립", C_MUTED),
]
for i, (name, desc, col) in enumerate(after):
    y = 2.0 + i * 1.1
    add_rect(sl, 7.2, y, 5.3, 0.85, C_CANVAS)
    add_rect(sl, 7.2, y, 0.06, 0.85, col)
    add_text(sl, name, 7.45, y + 0.07, 4.8, 0.32, size=12, bold=True, color=C_INK)
    add_text(sl, desc, 7.45, y + 0.42, 4.8, 0.32, size=11, color=C_MUTED)

add_text(sl, "→", 6.3, 3.8, 0.5, 0.5, size=28, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_text(sl, "비용 66% 절감", 6.1, 4.4, 0.9, 0.4, size=10, bold=True,
         color=C_SUCCESS, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 6 — EvidencePolicyVerifier 상세
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "EvidencePolicyVerifier — 증거 기반 취약점 검증")

add_text(sl, "Claude가 생성한 finding을 LLM 없이 결정론적 규칙으로 검증합니다.",
         0.7, 1.15, 12, 0.45, size=14, color=C_INK)

# 기본 체크
add_rect(sl, 0.5, 1.7, 6.0, 3.5, C_CARD)
add_rect(sl, 0.5, 1.7, 6.0, 0.06, C_PRIMARY)
add_text(sl, "기본 체크 (모든 severity)", 0.7, 1.78, 5.6, 0.4,
         size=13, bold=True, color=C_PRIMARY)

basic = [
    "evidence_quote 존재 여부",
    "evidence 소스 내 실제 존재 (Grounding)",
    "root_cause / remediation 필드 존재",
    "line_start 소스 라인 범위 이내",
]
for i, t in enumerate(basic):
    add_text(sl, f"✓  {t}", 0.7, 2.3 + i * 0.6, 5.6, 0.5, size=12, color=C_INK)

# Grounding 알고리즘
add_rect(sl, 0.5, 5.3, 6.0, 1.9, C_DARK)
add_text(sl, "Evidence Grounding 알고리즘", 0.7, 5.38, 5.6, 0.4,
         size=11, bold=True, color=C_ACCENT)
algo = [
    "① '...' / '…' 기준으로 fragment 분리",
    "② 각 fragment (≥8자)가 소스에 있으면 통과",
    "③ 12자 sliding window 매칭으로 폴백",
]
for i, t in enumerate(algo):
    add_text(sl, t, 0.7, 5.85 + i * 0.38, 5.6, 0.35, size=11, color=C_WHITE)

# Exploitability policy
add_rect(sl, 7.0, 1.7, 5.8, 5.5, C_CARD)
add_rect(sl, 7.0, 1.7, 5.8, 0.06, C_ERROR)
add_text(sl, "Exploitability Policy (severity별)", 7.2, 1.78, 5.4, 0.4,
         size=13, bold=True, color=C_ERROR)

rows = [
    ("CRITICAL", "≥ 70%", "VULNERABLE", "dangerous op 필수", C_ERROR),
    ("HIGH",     "≥ 60%", "VULNERABLE", "dangerous op 필수", C_AMBER),
    ("MEDIUM",   "≥ 50%", "VULNERABLE", "일반 기준",          C_WARNING),
    ("LOW",      "≥ 40%", "VUL / REVIEW", "일반 기준",        C_TEAL),
]
headers = ["Severity", "Confidence", "Verdict", "증거 요건"]
for j, h in enumerate(headers):
    xs = [7.2, 8.6, 9.8, 11.0]
    add_text(sl, h, xs[j], 2.3, 1.1, 0.35, size=10, bold=True, color=C_MUTED)

for i, (sev, conf, verdict, req, col) in enumerate(rows):
    y = 2.75 + i * 0.75
    add_rect(sl, 7.1, y, 5.5, 0.62, C_CANVAS)
    add_rect(sl, 7.1, y, 0.06, 0.62, col)
    add_text(sl, sev,    7.25, y + 0.15, 1.2, 0.35, size=11, bold=True, color=col)
    add_text(sl, conf,   8.55, y + 0.15, 1.1, 0.35, size=11, color=C_INK)
    add_text(sl, verdict,9.75, y + 0.15, 1.2, 0.35, size=11, color=C_INK)
    add_text(sl, req,    10.95, y + 0.15, 1.8, 0.35, size=10, color=C_MUTED)

add_text(sl, "* dangerous op: gets, strcpy, sprintf, memcpy, malloc(산술), scanf...",
         7.1, 6.6, 5.8, 0.4, size=9, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 7 — 프론트엔드 UI 구성
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "프론트엔드 — 디자인 시스템 & 주요 페이지")

pages = [
    ("HeroPage",          "분할 레이아웃\n터미널 목업 / 파이프라인 소개 / CTA"),
    ("ProjectsPage",      "프로젝트 목록\n생성 / 삭제 / 분석 현황 확인"),
    ("ProjectDetailPage", "ZIP 업로드\n분석 시작 / 리포트 히스토리"),
    ("ProgressPage",      "실시간 진행 현황\n5단계 스텝 + 로그 패널"),
    ("ReportPage",        "분석 결과 뷰어\nFinding 카드 / severity 분류 / PDF 다운로드"),
]
for i, (name, desc) in enumerate(pages):
    row, col = divmod(i, 3)
    x = 0.5 + col * 4.25
    y = 1.4 + row * 2.8
    add_rect(sl, x, y, 3.9, 2.4, C_CARD)
    add_rect(sl, x, y, 3.9, 0.06, C_PRIMARY)
    add_text(sl, name, x + 0.15, y + 0.12, 3.6, 0.4,
             size=13, bold=True, color=C_PRIMARY)
    add_text(sl, desc, x + 0.15, y + 0.62, 3.6, 1.6,
             size=11, color=C_INK)

# 디자인 토큰 요약
add_rect(sl, 0.5, 6.55, 12.3, 0.7, C_DARK)
tokens = [
    ("Canvas", "#faf9f5", C_CANVAS), ("Primary", "#014029", C_PRIMARY),
    ("Error",  "#c64545", C_ERROR),  ("Amber",   "#e8a55a", C_AMBER),
    ("Warning","#d4a017", C_WARNING),("Teal",    "#5db8a6", C_TEAL),
]
for i, (name, hex_, col) in enumerate(tokens):
    x = 0.8 + i * 2.0
    add_rect(sl, x, 6.65, 0.28, 0.28, col)
    add_text(sl, f"{name}  {hex_}", x + 0.35, 6.65, 1.55, 0.28,
             size=9, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 8 — PDF 리포트 구조
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "PDF 리포트 — 자동 생성 출력물")

# 페이지 미리보기 (텍스트로 묘사)
add_rect(sl, 0.5, 1.3, 4.8, 5.7, C_WHITE)
add_rect(sl, 0.5, 1.3, 4.8, 0.06, C_HAIRLINE)

# PDF 내부 구조 묘사
add_rect(sl, 0.6, 1.4, 4.6, 0.85, C_PRIMARY)
add_text(sl, "Reporter", 0.8, 1.45, 3.0, 0.45, size=18, bold=True, color=C_WHITE)
add_text(sl, "Security Analysis Report", 0.8, 1.82, 3.5, 0.3, size=10, color=C_ACCENT)
add_rect(sl, 3.8, 1.52, 1.1, 0.3, C_SUCCESS)
add_text(sl, "PASS", 3.8, 1.52, 1.1, 0.3, size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 메타정보
for i, line in enumerate(["Archive: openssl_1.0.2.zip", "Files analyzed: 20 / 312", "Accepted: 21  Rejected: 19"]):
    add_text(sl, line, 0.65, 2.35 + i*0.3, 4.5, 0.28, size=8, color=C_INK)

# 통계 박스
for i, (label, val, col) in enumerate([("FILES","20",C_PRIMARY),("ACCEPTED","21",C_SUCCESS),
                                        ("REJECTED","19",C_ERROR),("SKIPPED","2",C_MUTED)]):
    bx = 0.65 + i * 1.1
    add_rect(sl, bx, 3.35, 1.0, 0.6, C_CARD)
    add_rect(sl, bx, 3.35, 1.0, 0.05, col)
    add_text(sl, val, bx, 3.42, 1.0, 0.3, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, label, bx, 3.72, 1.0, 0.2, size=7, color=C_MUTED, align=PP_ALIGN.CENTER)

# Finding 블록
add_rect(sl, 0.65, 4.1, 4.5, 0.06, C_ERROR)
add_rect(sl, 0.65, 4.16, 4.5, 1.15, C_CARD)
add_text(sl, "buffer overflow in SSL_write", 0.8, 4.2, 4.2, 0.35, size=10, bold=True, color=C_INK)
add_text(sl, "CRITICAL  ·  L112  ·  94% confidence", 0.8, 4.55, 4.2, 0.28, size=8, color=C_ERROR)
add_rect(sl, 0.8, 4.9, 4.3, 0.35, C_DARK)
add_text(sl, '  strcpy(buf, user_data);  // no check', 0.85, 4.92, 4.2, 0.28,
         size=7, color=C_ACCENT)

# 우측 설명
features = [
    ("🎨  디자인", "웹 UI와 동일한 크림/다크그린 팔레트\nDejaVu Sans/Serif/Mono 폰트 임베딩"),
    ("📋  구조",   "헤더 → 메타 → 요약 → 통계 4칸\n→ 파일별 finding 노트"),
    ("💻  코드 블록", "evidence_quote를 별도 코드 박스\nseverity 색상 left stripe"),
    ("↓  자동 생성", "분석 완료 시 report.pdf 즉시 생성\n브라우저에서 직접 다운로드"),
]
for i, (title, desc) in enumerate(features):
    y = 1.35 + i * 1.5
    add_rect(sl, 6.0, y, 6.8, 1.3, C_CARD)
    add_rect(sl, 6.0, y, 6.8, 0.06, C_PRIMARY)
    add_text(sl, title, 6.2, y + 0.1, 6.4, 0.4, size=13, bold=True, color=C_INK)
    add_text(sl, desc, 6.2, y + 0.55, 6.4, 0.65, size=12, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 9 — AI 도구 활용 (Claude Code)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "AI 도구 활용 — Claude Code를 팀원처럼 매니징하기")

add_text(sl, "단순 자동완성이 아닌 '시니어 개발자 + AI'의 협업 구조로 프로젝트를 진행했습니다.",
         0.7, 1.15, 12, 0.45, size=14, color=C_INK)

cases = [
    ("아키텍처 설계 지시",
     "\"파일당 Claude API 3회 → 1회로 줄이고 싶다.\"\n→ Triage/Verify를 heuristic으로 교체하는 설계 제안 + 구현",
     C_PRIMARY),
    ("정책 결정 후 구현 위임",
     "\"Accepted 기준이 느슨하다. 실제 exploit 가능성을 더 봐야한다.\"\n→ severity별 confidence 임계값 + dangerous op 체크 설계 + 코드 작성",
     C_AMBER),
    ("디자인 시스템 방향 제시",
     "\"디자인을 design.md처럼 Anthropic 스타일로 바꿔줘.\"\n→ CSS 토큰, 컴포넌트, 레이아웃 전면 재작성",
     C_SUCCESS),
    ("버그 진단 및 수정",
     "\"evidence가 소스에 있는데 reject된다.\"\n→ grounding 알고리즘 분석 후 ellipsis 분리 + sliding window 방식으로 개선",
     C_TEAL),
    ("문서화 자동화",
     "\"지금까지 한 걸 가지고 docs에 md 만들어줘, README도.\"\n→ 구현 기록 + 아키텍처 문서 자동 생성 후 git push",
     C_MUTED),
]
for i, (title, desc, col) in enumerate(cases):
    row, c = divmod(i, 3)
    x = 0.5 + c * 4.25
    y = 1.75 + row * 2.5
    add_rect(sl, x, y, 3.9, 2.2, C_CARD)
    add_rect(sl, x, y, 0.07, 2.2, col)
    add_text(sl, title, x + 0.2, y + 0.1, 3.5, 0.4, size=12, bold=True, color=C_INK)
    add_text(sl, desc, x + 0.2, y + 0.58, 3.5, 1.5, size=10, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 10 — Git 커밋 히스토리
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "Git 활용 — 커밋 히스토리 & 버전 관리")

add_text(sl, "github.com/rohkyoungmin/AISECApp",
         0.7, 1.15, 12, 0.4, size=13, color=C_MUTED)

commits = [
    ("c71083e", "Add README and implementation docs",            "2025-05", C_MUTED),
    ("0aa6327", "Redesign UI, tighten verifier, improve PDF export", "2025-05", C_PRIMARY),
    ("f931630", "Add frontend project workflow",                 "2025-05", C_PRIMARY),
    ("27ab8cd", "Update default Claude model",                   "2025-05", C_SUCCESS),
    ("4a4d624", "Fix editable LLM install",                      "2025-04", C_SUCCESS),
    ("e875b22", "Use input directory for ZIP experiments",        "2025-04", C_AMBER),
    ("d1ec886", "Add experiment report exports",                  "2025-04", C_AMBER),
    ("f95c778", "Refactor source analysis into multi-agent pipeline", "2025-04", C_AMBER),
    ("4561353", "Add ZIP source analysis backend",               "2025-04", C_TEAL),
    ("2f0ea2f", "Add Claude source analysis report scaffold",     "2025-04", C_TEAL),
    ("0251ab7", "Import Magma patch case skeletons",             "2025-03", C_MUTED),
    ("e323228", "Initialize AISEC pipeline scaffold",            "2025-03", C_MUTED),
]
for i, (hash_, msg, date, col) in enumerate(commits):
    y = 1.65 + i * 0.47
    add_rect(sl, 0.5, y, 12.3, 0.42, C_CARD if i % 2 == 0 else C_CANVAS)
    add_rect(sl, 0.5, y, 0.06, 0.42, col)
    add_text(sl, hash_, 0.7, y + 0.07, 1.1, 0.3, size=9, color=C_MUTED)
    add_text(sl, msg,   1.85, y + 0.07, 9.5, 0.3, size=11, color=C_INK)
    add_text(sl, date,  11.5, y + 0.07, 1.2, 0.3, size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 11 — 트러블슈팅
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "트러블슈팅 — 기술적 난제와 해결 과정")

troubles = [
    (
        "① WSL2 환경 npm 충돌",
        "Windows용 npm이 PATH 우선순위를 가져 UNC 경로 오류\n(\\\\wsl.localhost\\... 경로를 처리 못함)",
        "nvm으로 Linux용 Node.js 20 LTS 별도 설치\nWindows node_modules 삭제 후 Linux npm으로 재설치",
        C_ERROR,
    ),
    (
        "② Evidence Grounding 과도한 Reject",
        "Claude가 '...'로 코드 일부 생략하거나 공백 다르게 출력\n→ 단순 substring 매칭으로 실제 존재하는 증거도 reject",
        "ellipsis 기준 fragment 분리 + 12자 sliding window 알고리즘으로 교체\n→ 정상 증거 reject 대폭 감소",
        C_AMBER,
    ),
    (
        "③ Exploitability 기준 모호",
        "evidence에 코드가 있어도 실제 exploit 가능 여부 판단 없음\n→ 이론적 취약점도 Accepted 처리됨",
        "severity별 confidence 임계값 + dangerous op 필수 체크 도입\nHIGH/CRITICAL은 gets/strcpy/memcpy 등 실제 위험 함수 증거 필수",
        C_WARNING,
    ),
    (
        "④ uvicorn --reload 미적용",
        "코드 변경 후 API 응답이 이전 코드 기준으로 동작\n→ --reload 없이 실행 중이었음",
        "--reload 옵션 추가 후 재시작\nStatReload가 소스 변경 자동 감지",
        C_TEAL,
    ),
]
for i, (title, prob, sol, col) in enumerate(troubles):
    row, c = divmod(i, 2)
    x = 0.5 + c * 6.5
    y = 1.3 + row * 2.9
    add_rect(sl, x, y, 6.2, 2.65, C_CARD)
    add_rect(sl, x, y, 6.2, 0.06, col)
    add_text(sl, title, x + 0.15, y + 0.1, 5.9, 0.4, size=13, bold=True, color=col)
    add_text(sl, "🔴 문제", x + 0.15, y + 0.58, 0.8, 0.3, size=10, bold=True, color=C_ERROR)
    add_text(sl, prob,  x + 0.95, y + 0.58, 5.1, 0.7, size=10, color=C_INK)
    add_text(sl, "✅ 해결", x + 0.15, y + 1.45, 0.8, 0.3, size=10, bold=True, color=C_SUCCESS)
    add_text(sl, sol,   x + 0.95, y + 1.45, 5.1, 0.85, size=10, color=C_INK)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 12 — 최종 데모까지의 계획
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
content_slide(sl, "최종 데모까지의 계획 — 5월 30일 Demo Day")

milestones = [
    ("5월 12일",  "중간 발표 자료 제출",           "완료",   C_SUCCESS),
    ("5월 13·18일","중간 발표 진행",              "예정",   C_PRIMARY),
    ("5월 20일",  "휴강",                        "-",      C_MUTED),
    ("5월 21~25일","CVE 케이스셋 정량 평가\nPrecision / Recall / F1 측정",
                  "진행 예정", C_AMBER),
    ("5월 26~28일","실시간 로그 스트리밍 (SSE)\nFalse Positive 분석 및 보정",
                  "진행 예정", C_AMBER),
    ("5월 29일",  "최종 시연 리허설 및 버그 수정",  "진행 예정", C_WARNING),
    ("5월 30일",  "최종 발표회 (Demo Day)\n13:00 ~ 16:00",
                  "목표",     C_ERROR),
]
for i, (date, task, status, col) in enumerate(milestones):
    y = 1.35 + i * 0.83
    add_rect(sl, 0.5, y, 12.3, 0.75, C_CARD if i % 2 == 0 else C_CANVAS)
    add_rect(sl, 0.5, y, 0.07, 0.75, col)
    add_text(sl, date,   0.7,  y + 0.18, 1.8, 0.4, size=12, bold=True, color=col)
    add_text(sl, task,   2.65, y + 0.05, 8.0, 0.65, size=12, color=C_INK)
    add_rect(sl, 11.2, y + 0.2, 1.4, 0.35, col)
    add_text(sl, status, 11.2, y + 0.2, 1.4, 0.35,
             size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 남은 To-do
add_rect(sl, 0.5, 7.05, 12.3, 0.25, C_PRIMARY)
add_text(sl, "핵심 남은 과제:  CVE 정량 평가  ·  SSE 실시간 스트리밍  ·  False Positive 보정  ·  Demo 리허설",
         0.7, 7.07, 12, 0.22, size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# 슬라이드 13 — 마무리
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, C_PRIMARY)
add_rect(sl, 0, 0, 13.33, 0.12, C_SUCCESS)

add_text(sl, "감사합니다", 0.7, 1.5, 12, 1.5,
         size=64, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Reporter — AI 기반 C/C++ 취약점 분석 플랫폼", 0.7, 3.2, 12, 0.7,
         size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 요약 3가지
for i, (icon, label) in enumerate([
    ("🤖", "Claude AI 멀티 에이전트\n파이프라인"),
    ("🔍", "증거 기반 검증\nExploitability Policy"),
    ("📄", "자동 PDF 리포트\n웹 UI 통합"),
]):
    x = 2.2 + i * 3.3
    add_rect(sl, x, 4.1, 2.9, 1.8, RGBColor(0x02, 0x55, 0x35))
    add_text(sl, icon,  x + 0.15, 4.2, 2.6, 0.5,
             size=24, align=PP_ALIGN.CENTER, color=C_WHITE)
    add_text(sl, label, x + 0.15, 4.8, 2.6, 0.9,
             size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_text(sl, "github.com/rohkyoungmin/AISECApp", 0.7, 6.4, 12, 0.5,
         size=13, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── 저장 ─────────────────────────────────────────────────────────────────────
out = "/home/rosette98/Projects/AISECApp/docs/Reporter_중간발표.pptx"
prs.save(out)
print(f"Saved: {out}")
